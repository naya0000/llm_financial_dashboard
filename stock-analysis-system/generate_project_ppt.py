from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path


OUTPUT = Path("/Users/naya/stock-analysis-system/stock-analysis-system-project-briefing-zh.pptx")


def set_text_style(text_frame, font_size=20, bold=False):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = "PingFang TC"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    set_text_style(slide.shapes.title.text_frame, 40, True)
    set_text_style(slide.placeholders[1].text_frame, 22, False)


def add_bullets_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    set_text_style(slide.shapes.title.text_frame, 34, True)

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.level = 0
        set_text_style(tf, 18, True)

    for i, bullet in enumerate(bullets):
        if i == 0 and not subtitle:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT

    set_text_style(tf, 20, False)


def build_presentation():
    prs = Presentation()

    add_title_slide(
        prs,
        "個人理財助理專案報告",
        "AI 多代理人個股分析系統\n專案分析與創新提案（10 分鐘）",
    )

    add_bullets_slide(
        prs,
        "目錄與報告節奏",
        [
            "1. 系統架構設計（2 分鐘）",
            "2. Prompt 設計：摘要精準與避免幻覺（2.5 分鐘）",
            "3. 資料品質控制：過濾不可靠/過時資訊（2 分鐘）",
            "4. 用戶體驗設計：回應簡潔、易懂（1.5 分鐘）",
            "5. 創新功能構想與落地路線（2 分鐘）",
        ],
    )

    add_bullets_slide(
        prs,
        "專案定位與目標",
        [
            "定位：Claude Code Skills 驅動的多代理人投資分析插件",
            "核心價值：將資料抓取、驗證、分析、整合、視覺化串成一條可靠流水線",
            "支援市場：台股、美股、日股、港股",
            "輸出型態：快速問答 + 結構化分析 + 互動式 Dashboard",
        ],
    )

    add_bullets_slide(
        prs,
        "系統架構總覽（端到端）",
        [
            "使用者輸入自然語言問題（例如：分析台積電）",
            "Orchestrator 解析意圖與模式（quick / selective / full）",
            "Fetcher + Validator 先建立可信資料基底",
            "6 位分析師並行輸出專業 JSON（財務、技術、量化、產業、情緒、法人）",
            "整合層產生 integrated_report，再由 Dashboard 腳本輸出 HTML 報告",
        ],
    )

    add_bullets_slide(
        prs,
        "系統架構設計亮點",
        [
            "混合架構：Python 負責機械穩定性，LLM 負責語義推理",
            "並行化：資料抓取 + 多分析師並行，降低整體等待時間",
            "快取策略：同日報告與 Agent 級快取，避免重複計算",
            "模式化服務：快速問答優先低延遲、完整分析優先深度",
        ],
    )

    add_bullets_slide(
        prs,
        "Prompt 設計原則（摘要精準）",
        [
            "流程強約束：明確 Step-by-Step，避免任意跳步與遺漏",
            "資料切片：每位 Agent 只吃自己需要的欄位，降低上下文噪音",
            "數值引用規範：引用指標必須給出精確值，不可籠統描述",
            "輸出格式固定：JSON schema + 指定欄位，降低摘要漂移",
        ],
    )

    add_bullets_slide(
        prs,
        "Prompt 設計原則（避免幻覺）",
        [
            "全域零幻覺政策：禁止用訓練資料補缺口",
            "缺資料必揭露：data_limitations + 摘要末段「⚠ 資料限制」",
            "信心度聯動：資料越缺，confidence 必須越低",
            "寧可留白不可捏造：無法驗證即明確說明無法評估",
        ],
    )

    add_bullets_slide(
        prs,
        "摘要精準機制：從 Agent 到整合層",
        [
            "Score-then-Justify：先初評再論證，避免事後硬湊敘事",
            "分數漂移審計：preliminary_score 與 final score 偏差過大需解釋",
            "限制整併規則：整合時必須保留所有關鍵限制，不可美化",
            "敘事分工：LLM 只做必要推理，機械欄位交給 Python 組裝",
        ],
    )

    add_bullets_slide(
        prs,
        "資料品質控制框架",
        [
            "三層品質閘門：hard_stop / warning / passed",
            "新鮮度檢查：價格 3 天內、財報 120 天內、新聞 30 天內",
            "異常偵測：PE/PB 極值、單日劇烈波動、量能異常",
            "市場化閾值：台股漲跌幅採 11% 邏輯，降低誤報",
        ],
    )

    add_bullets_slide(
        prs,
        "資料品質控制：決策行為",
        [
            "hard_stop（<30）：直接中止分析，避免錯誤結論輸出",
            "warning（30-49）：允許繼續，但明確警示並降低信心",
            "passed（>=50）：進入完整分析流程",
            "目標：把不可靠資料擋在分析前，而不是事後補救",
        ],
    )

    add_bullets_slide(
        prs,
        "用戶體驗設計：回應簡潔、易懂",
        [
            "三種模式自動切換：快問快答 / 選擇性分析 / 完整分析",
            "快速問答只回必要資訊（1-2 句），避免資訊過載",
            "趨勢問題提供固定結構：走勢圖 + 摘要表 + 趨勢解讀",
            "繁體中文專業語言，兼顧可讀性與決策可用性",
        ],
    )

    add_bullets_slide(
        prs,
        "用戶體驗設計：報告可理解性",
        [
            "Dashboard 統一視覺語言：評分、雷達圖、分析師觀點分區呈現",
            "重點先行：先給結論與信心，再展開依據",
            "可追溯性：新聞來源、資料限制、風險因子可回查",
            "降低認知負擔：技術細節由系統吸收，使用者看到決策重點",
        ],
    )

    add_bullets_slide(
        prs,
        "創新構想 1：自動驗證回圈（Auto-Refine）",
        [
            "當摘要與資料限制衝突時，自動退回整合層重新生成",
            "新增『矛盾偵測器』：檢查高信心但高缺失的異常組合",
            "對外輸出前加一道一致性簽核，進一步降低幻覺風險",
            "預期效益：提升報告可信度與可審計性",
        ],
    )

    add_bullets_slide(
        prs,
        "創新構想 2：使用者個人化策略助理",
        [
            "加入風險偏好檔案（保守/平衡/積極）",
            "同一份分析可輸出三種建議版本（防守、中性、進攻）",
            "新增『我該看哪三件事』摘要卡，提高行動可執行性",
            "預期效益：從資訊工具升級為決策輔助工具",
        ],
    )

    add_bullets_slide(
        prs,
        "創新構想 3：投組與事件驅動監控",
        [
            "多檔股票組合觀察：相關性、集中度、共通風險",
            "重大事件觸發再分析（財報、法說、監管）",
            "用戶只看『變化』而非重複全文，提升效率",
            "預期效益：從單次分析走向持續追蹤服務",
        ],
    )

    add_bullets_slide(
        prs,
        "落地時程與量化 KPI",
        [
            "Phase 1（2 週）：矛盾偵測 + 重新整合回圈",
            "Phase 2（3 週）：個人化建議模板與偏好設定",
            "Phase 3（4 週）：投組監控與事件觸發",
            "KPI：摘要正確率、幻覺事件數、平均回應時間、用戶採納率",
        ],
    )

    add_bullets_slide(
        prs,
        "結論",
        [
            "本系統的核心競爭力是『可信任的多代理人分析流程』",
            "Prompt 設計不是只寫提示詞，而是把品質與誠實性變成硬規則",
            "資料品質閘門 + 信心降級 + 限制揭露，形成防幻覺閉環",
            "下一步可透過個人化與持續監控，擴展成長期投資決策平台",
        ],
    )

    add_bullets_slide(
        prs,
        "Q&A",
        [
            "謝謝聆聽",
            "如需，我可以再提供：",
            "- 主管版（8 分鐘）精簡稿",
            "- 技術評審版（15 分鐘）深度稿",
            "- 可直接演講的逐字講稿",
        ],
    )

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
